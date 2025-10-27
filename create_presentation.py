#!/usr/bin/env python3
"""
Professional Presentation Generator
신뢰감 있고 정성스러운 디자인의 PPTX 생성 도구
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json


class ProfessionalPresentation:
    """전문적인 프레젠테이션 생성 클래스"""

    # 신뢰감 있는 색상 팔레트 (블루 톤 - 안정감, 전문성)
    COLORS = {
        'primary': RGBColor(30, 58, 95),      # 진한 청색 #1E3A5F
        'secondary': RGBColor(52, 152, 219),   # 밝은 청색 #3498DB
        'accent': RGBColor(243, 156, 18),      # 골드 #F39C12
        'dark_gray': RGBColor(52, 73, 94),     # 진한 회색 #34495E
        'light_gray': RGBColor(236, 240, 241), # 연한 회색 #ECF0F1
        'white': RGBColor(255, 255, 255),      # 흰색
        'text': RGBColor(44, 62, 80),          # 텍스트 #2C3E50
    }

    # 폰트 크기
    FONT_SIZES = {
        'title': Pt(44),
        'subtitle': Pt(24),
        'heading': Pt(32),
        'subheading': Pt(24),
        'body': Pt(18),
        'caption': Pt(14),
    }

    def __init__(self):
        """프레젠테이션 초기화"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle, author=""):
        """
        타이틀 슬라이드 추가

        Args:
            title: 메인 제목
            subtitle: 부제목
            author: 작성자 (선택사항)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 빈 레이아웃

        # 배경 장식 추가 (상단 헤더)
        header = slide.shapes.add_shape(
            1,  # 사각형
            Inches(0), Inches(0),
            Inches(10), Inches(1.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS['primary']
        header.line.fill.background()

        # 메인 제목
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = self.FONT_SIZES['title']
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']

        # 부제목
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(8), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = self.FONT_SIZES['subtitle']
        subtitle_para.font.color.rgb = self.COLORS['dark_gray']

        # 작성자 정보 (하단)
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5),
                Inches(8), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.text = author
            author_para = author_frame.paragraphs[0]
            author_para.alignment = PP_ALIGN.CENTER
            author_para.font.size = self.FONT_SIZES['body']
            author_para.font.color.rgb = self.COLORS['dark_gray']

        # 하단 장식 라인
        footer_line = slide.shapes.add_shape(
            1,  # 사각형
            Inches(3), Inches(6.8),
            Inches(4), Inches(0.05)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.COLORS['accent']
        footer_line.line.fill.background()

        return slide

    def add_section_slide(self, section_title):
        """
        섹션 구분 슬라이드 추가

        Args:
            section_title: 섹션 제목
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 배경 그라데이션 효과 (사각형 3개로 표현)
        bg1 = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = self.COLORS['primary']
        bg1.line.fill.background()

        # 섹션 제목
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = self.FONT_SIZES['heading']
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']

        # 장식 라인
        accent_line = slide.shapes.add_shape(
            1, Inches(3.5), Inches(5),
            Inches(3), Inches(0.1)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.COLORS['accent']
        accent_line.line.fill.background()

        return slide

    def add_content_slide(self, title, content_items, layout='bullet'):
        """
        내용 슬라이드 추가

        Args:
            title: 슬라이드 제목
            content_items: 내용 리스트 또는 딕셔너리
            layout: 'bullet' 또는 'two_column'
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 헤더 영역
        self._add_header(slide, title)

        if layout == 'bullet':
            self._add_bullet_content(slide, content_items)
        elif layout == 'two_column':
            self._add_two_column_content(slide, content_items)

        return slide

    def add_image_slide(self, title, image_path, caption=""):
        """
        이미지 슬라이드 추가

        Args:
            title: 슬라이드 제목
            image_path: 이미지 파일 경로
            caption: 이미지 설명 (선택사항)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 헤더 영역
        self._add_header(slide, title)

        # 이미지 추가 (중앙에 배치)
        if Path(image_path).exists():
            try:
                # 이미지 크기 조정 (슬라이드 크기의 70%)
                img = slide.shapes.add_picture(
                    image_path,
                    Inches(1.5), Inches(2),
                    width=Inches(7)
                )
            except Exception as e:
                # 이미지 로드 실패 시 placeholder 추가
                self._add_image_placeholder(slide)
        else:
            self._add_image_placeholder(slide)

        # 캡션 추가
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6),
                Inches(7), Inches(0.8)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_para = caption_frame.paragraphs[0]
            caption_para.alignment = PP_ALIGN.CENTER
            caption_para.font.size = self.FONT_SIZES['caption']
            caption_para.font.italic = True
            caption_para.font.color.rgb = self.COLORS['dark_gray']

        return slide

    def add_conclusion_slide(self, title, key_points):
        """
        결론 슬라이드 추가

        Args:
            title: 결론 제목
            key_points: 핵심 포인트 리스트 (최대 3개 권장)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 배경 장식
        bg_accent = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(3), Inches(7.5)
        )
        bg_accent.fill.solid()
        bg_accent.fill.fore_color.rgb = self.COLORS['primary']
        bg_accent.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(3.5), Inches(1),
            Inches(6), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = self.FONT_SIZES['heading']
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['primary']

        # 핵심 포인트
        y_position = 2.5
        for i, point in enumerate(key_points[:3]):  # 최대 3개
            # 번호 박스
            num_box = slide.shapes.add_shape(
                1, Inches(3.5), Inches(y_position),
                Inches(0.6), Inches(0.6)
            )
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = self.COLORS['accent']
            num_box.line.fill.background()

            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_para = num_frame.paragraphs[0]
            num_para.alignment = PP_ALIGN.CENTER
            num_para.font.size = self.FONT_SIZES['subheading']
            num_para.font.bold = True
            num_para.font.color.rgb = self.COLORS['white']

            # 포인트 텍스트
            point_box = slide.shapes.add_textbox(
                Inches(4.3), Inches(y_position),
                Inches(5.2), Inches(0.8)
            )
            point_frame = point_box.text_frame
            point_frame.text = point
            point_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            point_para = point_frame.paragraphs[0]
            point_para.font.size = self.FONT_SIZES['body']
            point_para.font.color.rgb = self.COLORS['text']

            y_position += 1.2

        return slide

    def _add_header(self, slide, title):
        """헤더 영역 추가"""
        # 헤더 배경
        header_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.COLORS['primary']
        header_bg.line.fill.background()

        # 제목 텍스트
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_para = title_frame.paragraphs[0]
        title_para.font.size = self.FONT_SIZES['subheading']
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']

        # 하단 액센트 라인
        accent_line = slide.shapes.add_shape(
            1, Inches(0), Inches(1.15),
            Inches(10), Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.COLORS['accent']
        accent_line.line.fill.background()

    def _add_bullet_content(self, slide, items):
        """불릿 포인트 내용 추가"""
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(items):
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = "• " + item
            p.font.size = self.FONT_SIZES['body']
            p.font.color.rgb = self.COLORS['text']
            p.space_before = Pt(12)
            p.level = 0

    def _add_two_column_content(self, slide, content_dict):
        """2단 컬럼 내용 추가"""
        # 왼쪽 컬럼
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8),
            Inches(4), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        if 'left' in content_dict:
            for i, item in enumerate(content_dict['left']):
                if i > 0:
                    left_frame.add_paragraph()
                p = left_frame.paragraphs[i]
                p.text = "• " + item
                p.font.size = self.FONT_SIZES['body']
                p.font.color.rgb = self.COLORS['text']
                p.space_before = Pt(10)

        # 오른쪽 컬럼
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.8),
            Inches(4), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        if 'right' in content_dict:
            for i, item in enumerate(content_dict['right']):
                if i > 0:
                    right_frame.add_paragraph()
                p = right_frame.paragraphs[i]
                p.text = "• " + item
                p.font.size = self.FONT_SIZES['body']
                p.font.color.rgb = self.COLORS['text']
                p.space_before = Pt(10)

        # 중앙 구분선
        divider = slide.shapes.add_shape(
            1, Inches(4.95), Inches(1.8),
            Inches(0.05), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.COLORS['light_gray']
        divider.line.fill.background()

    def _add_image_placeholder(self, slide):
        """이미지 placeholder 추가"""
        placeholder = slide.shapes.add_shape(
            1, Inches(2), Inches(2.5),
            Inches(6), Inches(3.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = self.COLORS['light_gray']
        placeholder.line.color.rgb = self.COLORS['dark_gray']

        # "Image" 텍스트
        text_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.8),
            Inches(6), Inches(0.8)
        )
        text_frame = text_box.text_frame
        text_frame.text = "[ 이미지 ]"
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = self.FONT_SIZES['heading']
        para.font.color.rgb = self.COLORS['dark_gray']

    def save(self, filename):
        """프레젠테이션 저장"""
        self.prs.save(filename)
        print(f"✓ 프레젠테이션이 저장되었습니다: {filename}")


def create_sample_presentation():
    """샘플 프레젠테이션 생성 (사용 예시)"""
    ppt = ProfessionalPresentation()

    # 1. 타이틀 슬라이드
    ppt.add_title_slide(
        title="프로젝트 보고서",
        subtitle="핵심 내용 요약 및 분석",
        author="작성자: Your Name | 날짜: 2025-10-27"
    )

    # 2. 목차 슬라이드
    ppt.add_content_slide(
        title="목차",
        content_items=[
            "1. 프로젝트 개요",
            "2. 주요 성과",
            "3. 기술적 분석",
            "4. 향후 계획",
            "5. 결론"
        ]
    )

    # 3. 섹션 구분
    ppt.add_section_slide("01. 프로젝트 개요")

    # 4. 내용 슬라이드
    ppt.add_content_slide(
        title="프로젝트 목표 및 범위",
        content_items=[
            "블록체인 시스템 컨트랙트 구현",
            "거버넌스 기능 통합 및 테스트",
            "성능 최적화 및 보안 강화",
            "완전한 테스트 커버리지 달성"
        ]
    )

    # 5. 2단 컬럼 슬라이드
    ppt.add_content_slide(
        title="주요 특징",
        content_items={
            'left': [
                "모듈화된 아키텍처",
                "확장 가능한 설계",
                "완전한 문서화"
            ],
            'right': [
                "높은 테스트 커버리지",
                "보안 강화",
                "성능 최적화"
            ]
        },
        layout='two_column'
    )

    # 6. 섹션 구분
    ppt.add_section_slide("02. 주요 성과")

    # 7. 이미지 슬라이드 예시
    ppt.add_image_slide(
        title="시스템 아키텍처",
        image_path="architecture_diagram.png",
        caption="그림 1. 전체 시스템 구조도"
    )

    # 8. 결론
    ppt.add_conclusion_slide(
        title="핵심 요약",
        key_points=[
            "성공적인 시스템 통합 및 테스트 완료",
            "높은 품질과 보안성 확보",
            "확장 가능한 아키텍처 구현"
        ]
    )

    # 저장
    ppt.save("professional_presentation.pptx")


def create_from_json(json_file):
    """JSON 파일에서 프레젠테이션 생성"""
    with open(json_file, 'r', encoding='utf-8') as f:
        data = json.load(f)

    ppt = ProfessionalPresentation()

    # 타이틀
    if 'title' in data:
        ppt.add_title_slide(
            title=data['title'].get('main', ''),
            subtitle=data['title'].get('sub', ''),
            author=data['title'].get('author', '')
        )

    # 슬라이드
    for slide_data in data.get('slides', []):
        slide_type = slide_data.get('type', 'content')

        if slide_type == 'section':
            ppt.add_section_slide(slide_data.get('title', ''))
        elif slide_type == 'content':
            ppt.add_content_slide(
                title=slide_data.get('title', ''),
                content_items=slide_data.get('content', []),
                layout=slide_data.get('layout', 'bullet')
            )
        elif slide_type == 'image':
            ppt.add_image_slide(
                title=slide_data.get('title', ''),
                image_path=slide_data.get('image', ''),
                caption=slide_data.get('caption', '')
            )
        elif slide_type == 'conclusion':
            ppt.add_conclusion_slide(
                title=slide_data.get('title', ''),
                key_points=slide_data.get('points', [])
            )

    # 저장
    output_file = data.get('output', 'presentation.pptx')
    ppt.save(output_file)


if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1:
        # JSON 파일에서 생성
        create_from_json(sys.argv[1])
    else:
        # 샘플 프레젠테이션 생성
        print("샘플 프레젠테이션을 생성합니다...")
        create_sample_presentation()
